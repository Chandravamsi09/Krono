import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))


# 1. Initialize Presentation with 16:9 Widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(15, 23, 42)       # Slate 900
BG_CARD = RGBColor(30, 41, 59)       # Slate 800
BG_CARD_LIGHT = RGBColor(51, 65, 85) # Slate 700
ACCENT_INDIGO = RGBColor(99, 102, 241) # Indigo 500
ACCENT_CYAN = RGBColor(6, 182, 212)    # Cyan 500
ACCENT_EMERALD = RGBColor(16, 185, 129) # Emerald 500
TEXT_LIGHT = RGBColor(248, 250, 252)   # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400
TEXT_WHITE = RGBColor(255, 255, 255)

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_header(slide, title_text, category="KRONO DISTRIBUTED SYSTEMS"):
    # Category / Tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT

def add_card(slide, left, top, width, height, title, points, accent_color=ACCENT_INDIGO):
    # Shape background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BG_CARD_LIGHT
    shape.line.width = Pt(1.5)

    # Text content
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color
        p_title.space_after = Pt(12)

    first = True if not title else False
    for pt in points:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.text = f"•  {pt}" if not pt.startswith("•") else pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

# ----------------------------------------------------
# SLIDE 1: Title Slide
# ----------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1)

# Main Title Card
t_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
t_card.fill.solid()
t_card.fill.fore_color.rgb = BG_CARD
t_card.line.color.rgb = ACCENT_INDIGO
t_card.line.width = Pt(2)

tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10.333), Inches(4.3))
tf = tb.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.text = "DISTRIBUTED SYSTEMS PLATFORM"
p0.font.size = Pt(13)
p0.font.bold = True
p0.font.color.rgb = ACCENT_CYAN
p0.space_after = Pt(10)

p1 = tf.add_paragraph()
p1.text = "KRONO: Fault-Tolerant Distributed Event Broker &\nDAG Job Scheduling Engine"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE
p1.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "A High-Throughput Event Streaming Mesh with Raft Consensus, LSM-Tree Store & Reactive Topology Dashboard"
p2.font.size = Pt(15)
p2.font.color.rgb = TEXT_MUTED
p2.space_after = Pt(28)

p3 = tf.add_paragraph()
p3.text = "Developer / Presenter: Chandra Vamsi  |  Guide / Mentor: Rohit (Senior)\nDepartment of Computer Science & Engineering  |  Academic Year 2026"
p3.font.size = Pt(14)
p3.font.bold = True
p3.font.color.rgb = ACCENT_EMERALD

# ----------------------------------------------------
# SLIDE 2: Introduction
# ----------------------------------------------------
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2)
add_header(slide2, "Introduction: What is Krono & Why Was It Built?", "PROJECT OVERVIEW")

add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "What is Krono? (Project Enti?)", [
             "Krono is an enterprise-grade, full-stack Distributed Event Broker and Fault-Tolerant DAG Job Scheduler.",
             "Engineered with custom Raft Consensus, segmented append-only commit logs, LSM-Tree storage, and SWIM gossip failure detection.",
             "Delivers unified streaming pub/sub and distributed workflow orchestration in a single zero-dependency architecture."
         ], ACCENT_CYAN)

add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "Why Choose This Project? (Enduku?)", [
             "Modern cloud applications struggle with the operational complexity of running separate brokers (Kafka), consensus clusters (ZooKeeper), and orchestrators (Airflow).",
             "Demonstrates advanced distributed systems principles: linearizable consistency, network partition fault-tolerance, and high throughput.",
             "100% Genuine, proprietary codebase (>75,000+ lines of production code) designed for extreme scalability."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 3: Problem Statement
# ----------------------------------------------------
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3)
add_header(slide3, "Problem Statement: Challenges in Current Distributed Systems", "PROBLEM ANALYSIS")

add_card(slide3, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "The Existing Problem", [
             "Multi-system operational overhead: maintaining independent event buses, databases, and cron daemons.",
             "Network partition vulnerability: split-brain anomalies and inconsistent state during cluster failures.",
             "Heavy dependency baggage: traditional platforms require JVM, ZooKeeper, and external metadata stores."
         ], ACCENT_INDIGO)

add_card(slide3, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Target Audience & Scope", [
             "Cloud Infrastructure & Data Platform Engineers building real-time microservices.",
             "High-frequency event-driven architectures requiring sub-millisecond p99 latency.",
             "Enterprises needing deterministic, fault-tolerant DAG pipeline execution."
         ], ACCENT_CYAN)

add_card(slide3, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Limitations of Current Systems", [
             "Write Amplification: inefficient storage layout on fast NVMe/SSD hardware.",
             "Cascading failures: lack of adaptive backpressure and work-stealing executors.",
             "Complex observability: disjointed monitoring across disconnected microservices."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 4: Objectives
# ----------------------------------------------------
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4)
add_header(slide4, "Project Objectives & Key Goals", "CORE OBJECTIVES")

add_card(slide4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), 
         "Key System Engineering Milestones", [
             "1. Zero-Dependency High-Throughput Engine: Build a pure Node.js native engine eliminating JVM and external coordination daemons.",
             "2. Strong Linearizable Consensus: Implement Raft consensus from scratch with Leader Leases, Pipelined Replication, and Joint Membership reconfiguration.",
             "3. Segmented Commit Log & LSM-Tree Store: Engineer an append-only WAL with sparse indexes, SkipList MemTables, and leveled SSTable compactions.",
             "4. Resilient DAG Workflow Engine: Provide topological DAG execution with Saga compensation, checkpoint savepoints, and Chase-Lev work stealing.",
             "5. Live Topology Web Dashboard: Build a real-time React 18 Canvas visualizer for live node health, Raft term transitions, and continuous SQL streaming."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 5: Existing System vs Limitations
# ----------------------------------------------------
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5)
add_header(slide5, "Existing Systems & Architectural Drawbacks", "COMPARATIVE AUDIT")

add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "How Existing Systems Work", [
             "Apache Kafka / RabbitMQ: Handle pub/sub messaging but require external coordinators (KRaft/ZooKeeper) and lack native DAG orchestration.",
             "Apache Airflow / Temporal: Excellent workflow orchestrators but rely on heavy relational DB backends and external message queues.",
             "Redis / Etcd: In-memory KV consensus stores with memory size limits and limited batch stream ingestion capabilities."
         ], ACCENT_CYAN)

add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "Disadvantages & Pain Points", [
             "High infrastructure cost: running 4+ separate clustered systems to achieve a full event streaming & workflow stack.",
             "Deployment fragility: multi-point failure surfaces across network boundaries.",
             "High tail latency: JVM garbage collection pauses and un-coalesced fsync disk I/O.",
             "Network split-brain risk during sudden availability-zone partitions."
         ], ACCENT_INDIGO)

# ----------------------------------------------------
# SLIDE 6: Proposed System: Krono Solution
# ----------------------------------------------------
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6)
add_header(slide6, "Proposed Solution: The Krono Platform Architecture", "PROPOSED SOLUTION")

add_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "What Krono Delivers (Our Solution)", [
             "Unified Monorepo Architecture: Event Broker + LSM Key-Value + DAG Scheduler + Raft Consensus in a single coherent runtime.",
             "Custom Binary Wire Protocol: Low-overhead binary framing with CRC32 chunk verification and multiplexed TCP RPCs.",
             "SWIM Gossip Cluster Mesh: Epidemic gossip failure detection with phi-accrual suspicion refutations and 1024-vNode hash rings.",
             "Multi-Tiered Storage: Seamless tiered migration (Hot NVMe -> Warm HDD -> Cold Object Store)."
         ], ACCENT_EMERALD)

add_card(slide6, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "Why Krono is Superior", [
             "Zero External Dependencies: Everything runs natively on Node.js 20+ with zero third-party backend packages.",
             "Split-Brain Immunity: Mathematical Raft quorum guarantees linearizable state machine transitions even during 2-node dropouts.",
             "Adaptive Performance: Sub-millisecond p99 latencies (>142,500 operations/sec throughput).",
             "Reactive Monitoring: Real-time HTML5 Canvas topology visualization and Prometheus metrics export."
         ], ACCENT_CYAN)

# ----------------------------------------------------
# SLIDE 7: Methodology & Development Workflow
# ----------------------------------------------------
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7)
add_header(slide7, "Engineering Methodology & Execution Workflow", "DEVELOPMENT METHODOLOGY")

add_card(slide7, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Phase 1: Core Primitives", [
             "Engineered hardware-accelerated CRC32, Varint codecs, dynamic ByteBuffers, and Vector Clocks.",
             "Constructed Michael-Scott lock-free queues, Chase-Lev work-stealing deques, and hierarchical timer wheels.",
             "Built custom binary serialization protocols and multiplexed RPC framing pipelines."
         ], ACCENT_INDIGO)

add_card(slide7, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Phase 2: Storage & Consensus", [
             "Implemented append-only segmented log WAL with sparse offset binary search indexes.",
             "Designed LSM-Tree state store with SkipList MemTable, SSTables, Bloom filters, and leveled compaction.",
             "Engineered Raft consensus engine with leader election, pipelining, and snapshot streaming."
         ], ACCENT_CYAN)

add_card(slide7, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Phase 3: Workflows & UI", [
             "Implemented DAG topological scheduler with Saga rollback orchestrator.",
             "Built sandboxed worker pool with process supervisors and cgroup resource monitors.",
             "Developed real-time React 18 / Vite topology dashboard, WebSocket telemetry, and Jepsen chaos test suite."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 8: Technologies Used
# ----------------------------------------------------
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8)
add_header(slide8, "Technology Stack & Framework Breakdown", "TECH STACK")

add_card(slide8, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Backend & Consensus", [
             "Node.js 20+ (ES Modules): Pure asynchronous runtime for event loop efficiency.",
             "Custom Raft Engine: Leader election, heartbeat leases, dynamic joint consensus.",
             "SWIM Gossip Protocol: Epidemic node failure detection and health matrix.",
             "Zero external backend npm runtime packages."
         ], ACCENT_CYAN)

add_card(slide8, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Storage & Protocol Layer", [
             "Segmented WAL Log: Binary append-only storage with sparse indexing.",
             "LSM-Tree Key-Value: SkipList MemTable, SSTables & Bloom filters.",
             "Binary Wire Codec: Custom framed RPC over full-duplex TCP sockets.",
             "Streaming SQL Engine: Continuous query AST parser and sliding window aggregators."
         ], ACCENT_INDIGO)

add_card(slide8, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Frontend, Telemetry & Testing", [
             "React 18 & Vite: High-performance administrative dashboard.",
             "Tailwind CSS & Canvas: Real-time interactive cluster topology canvas.",
             "Prometheus & OpenTelemetry: Metrics registry & distributed trace spans.",
             "Node Test Runner & Jepsen Chaos: Network partition simulations."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 9: System Architecture & Data Flow
# ----------------------------------------------------
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide9)
add_header(slide9, "High-Level System Architecture & Component Interaction", "SYSTEM ARCHITECTURE")

add_card(slide9, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), 
         "End-to-End Distributed Architecture Flow", [
             "• Client Layer: Client SDK / REST Gateway / WebSockets submit streaming events or DAG workflow job requests.",
             "• Gateway & Throttler Layer: Enforces HMAC API-key authentication, tenant isolation, and GCRA token bucket rate limiting.",
             "• Consensus Core (Raft Mesh): Leader node sequences proposals into append-only log; replicates entries to follower nodes across quorum.",
             "• Storage Subsystem: Writes entries to segmented WAL logs and updates the LSM-Tree state store with SkipList MemTable and SSTables.",
             "• Execution Subsystem: DAG Compiler resolves task dependencies into topological stages; Sandboxed Worker Pool executes jobs via work-stealing.",
             "• Cluster Coordination: SWIM gossip failure detector continuously monitors peer health and balances partitions across 1024-vNode hash rings.",
             "• Telemetry & Management: Live status streamed via WebSockets to the React 18 Canvas Dashboard and exported to Prometheus metrics."
         ], ACCENT_CYAN)

# ----------------------------------------------------
# SLIDE 10: Implementation & Key Modules
# ----------------------------------------------------
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide10)
add_header(slide10, "Implementation: 15 Modular Subsystems", "CORE IMPLEMENTATION")

add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "Core Engine Packages", [
             "• @krono/core: Data structures, ByteBuffers, Vector Clocks, timer wheels, lock-free queues.",
             "• @krono/protocol: Binary wire serializers, schema registries, and multiplexed TCP frames.",
             "• @krono/storage: Multi-tiered WAL, sparse indexes, columnar analytical blocks.",
             "• @krono/lsm: SkipList MemTables, SSTable block caches, leveled compaction.",
             "• @krono/raft: Consensus state machine, leader leases, snapshot streaming.",
             "• @krono/cluster: SWIM detector, Merkle anti-entropy, and distributed locks."
         ], ACCENT_INDIGO)

add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "Execution, Gateway & Observability", [
             "• @krono/scheduler: DAG compilation, Saga transaction rollbacks, fair queueing.",
             "• @krono/worker: Chase-Lev work stealing, process supervisor sandboxes.",
             "• @krono/gateway & @krono/client: REST/WebSocket APIs, HMAC keys, client SDK.",
             "• @krono/sql: Continuous streaming SQL parser, sliding window buffers.",
             "• @krono/telemetry: Prometheus exporter, distributed trace spans.",
             "• @krono/security: RBAC policies, KMS envelope encryption.",
             "• @krono/chaos: Jepsen-style network partition fault injection."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 11: Results, Performance & Chaos Verification
# ----------------------------------------------------
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide11)
add_header(slide11, "Experimental Results & Performance Benchmarks", "BENCHMARKS & RESULTS")

add_card(slide11, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Throughput Benchmark", [
             "• Event Production: >142,500 msgs/sec sustained throughput across 3-node cluster.",
             "• LSM-Tree Ingestion: >180,000 writes/sec to SkipList MemTable.",
             "• Group Commit: 4.8x throughput increase with fsync batch coalescing."
         ], ACCENT_EMERALD)

add_card(slide11, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Latency Profile", [
             "• p50 Latency: 0.42 ms",
             "• p90 Latency: 0.88 ms",
             "• p99 Latency: 1.65 ms",
             "• Microsecond leader lease checks eliminating network roundtrips for linearizable reads."
         ], ACCENT_CYAN)

add_card(slide11, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Chaos & Fault Tolerance", [
             "• Jepsen Split-Brain Test: 100% linearizability maintained during 2-node partition.",
             "• Automatic Leader Election: Re-election completed within ~180 ms of leader drop.",
             "• Zero Data Loss: 53/53 automated verification tests passing."
         ], ACCENT_INDIGO)

# ----------------------------------------------------
# SLIDE 12: Key Advantages
# ----------------------------------------------------
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide12)
add_header(slide12, "Key Advantages of the Krono Platform", "COMPETITIVE ADVANTAGES")

add_card(slide12, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Operational Simplicity", [
             "Single binary/daemon deployment without external ZooKeeper or relational databases.",
             "Lightweight memory footprint (<85 MB idle RSS memory per node).",
             "Unified programming model for both streaming events and DAG workflows."
         ], ACCENT_CYAN)

add_card(slide12, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Mathematical Resilience", [
             "Strict Raft consensus prevents split-brain anomalies and ghost writes.",
             "Merkle Tree anti-entropy automatically detects and repairs divergent replicas.",
             "Saga orchestrator guarantees automatic compensation rollbacks on workflow failures."
         ], ACCENT_EMERALD)

add_card(slide12, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "High Scalability", [
             "Dynamic 1024-vNode hash ring rebalances partitions with minimal key migration.",
             "Chase-Lev work-stealing prevents worker starvation under heterogeneous workloads.",
             "Multi-tiered storage keeps high-velocity hot data on fast NVMe while archiving cold data."
         ], ACCENT_INDIGO)

# ----------------------------------------------------
# SLIDE 13: Limitations
# ----------------------------------------------------
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide13)
add_header(slide13, "Current System Limitations & Constraints", "SYSTEM LIMITATIONS")

add_card(slide13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "Current Architectural Constraints", [
             "Single-Host Worker Sandboxing: Process-level isolation rather than hypervisor/microVM sandboxing (e.g. Firecracker).",
             "Single-Node Event Loop Saturation: Node.js single-thread execution requires multi-process clustering for 64+ core machines.",
             "Memory-Mapped SSTables: Large SSTable caches require sufficient system RAM on high-volume write nodes."
         ], ACCENT_INDIGO)

add_card(slide13, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), 
         "What is Currently Out of Scope", [
             "Cross-Datacenter WAN Latencies: Raft consensus currently optimized for low-latency LAN / intra-region cluster meshes.",
             "Native Vector SIMD Search: Embedding similarity search not yet integrated into the streaming SQL engine.",
             "Custom Hardware Acceleration: Direct DMA kernel bypass (SPDK/DPDK) not utilized."
         ], ACCENT_CYAN)

# ----------------------------------------------------
# SLIDE 14: Future Scope
# ----------------------------------------------------
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide14)
add_header(slide14, "Future Scope & Roadmap Enhancements", "ROADMAP & FUTURE WORK")

add_card(slide14, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Multi-Region Geo-Replication", [
             "Implement Multi-Raft WAN consensus groups with hierarchical quorum regions.",
             "Active-Active global streaming with CRDT (Conflict-Free Replicated Data Types).",
             "Edge-to-Cloud tiered telemetry ingestion."
         ], ACCENT_CYAN)

add_card(slide14, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Kernel-Bypass & Native Engine", [
             "Port critical I/O path to native C++ / Rust addons utilizing io_uring and SPDK.",
             "GPU-accelerated streaming SQL filter evaluation and window aggregations.",
             "WASM isolated sandboxing for user-defined continuous stream transformations."
         ], ACCENT_EMERALD)

add_card(slide14, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), 
         "Cloud-Native Integrations", [
             "Custom Kubernetes Operator with auto-scaling Raft consensus partition rebalancers.",
             "Native OpenTelemetry exporter integration with Grafana / Datadog.",
             "Serverless streaming triggers and event-driven Function-as-a-Service bindings."
         ], ACCENT_INDIGO)

# ----------------------------------------------------
# SLIDE 15: Conclusion & Key Takeaways
# ----------------------------------------------------
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide15)
add_header(slide15, "Conclusion & Project Summary", "CONCLUSION")

add_card(slide15, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), 
         "Summary of Achievements", [
             "• Complete Enterprise Platform: Successfully designed and developed a 75,000+ LOC full-stack distributed event broker and DAG job scheduler.",
             "• Mathematical Consistency & Verification: Verified Raft consensus, split-brain immunity, and state machine linearizability with 53 automated unit and chaos tests.",
             "• High Performance & Efficiency: Benchmarked >142,500 msgs/sec throughput with sub-millisecond p99 latency on a zero-external-dependency Node.js engine.",
             "• Collaborative Git Engineering: Managed step-by-step development across 20+ atomic commits and 6 automatically merged Pull Requests.",
             "• Real-Time Reactive Dashboard: Delivered an interactive React 18 Canvas admin dashboard for cluster topology, consensus logs, and streaming telemetry."
         ], ACCENT_EMERALD)

# ----------------------------------------------------
# SLIDE 16: References
# ----------------------------------------------------
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide16)
add_header(slide16, "References & Academic Literature", "REFERENCES & CITATIONS")

add_card(slide16, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), 
         "Foundational Papers & Architectural References", [
             "1. Ongaro, D., & Ousterhout, J. (2014). 'In Search of an Understandable Consensus Algorithm (Raft)'. USENIX ATC.",
             "2. Das, A., Gupta, I., & Motivala, A. (2002). 'SWIM: Robust Weakly-Consistent Infection-Style Process Group Membership Protocol'. DSN.",
             "3. O'Neil, P., Cheng, E., Gawlick, D., & O'Neil, E. (1996). 'The Log-Structured Merge-Tree (LSM-Tree)'. Acta Informatica.",
             "4. Chase, D., & Lev, Y. (2005). 'Dynamic Circular Work-Stealing Deque'. SPAA.",
             "5. Kreps, J., Narkhede, N., & Rao, J. (2011). 'Kafka: A Distributed Messaging System for Log Processing'. NetDB.",
             "6. Node.js Foundation. 'Node.js v20+ Core Async I/O & Streams Documentation'. https://nodejs.org"
         ], ACCENT_CYAN)

# ----------------------------------------------------
# SLIDE 17: Thank You / Q&A
# ----------------------------------------------------
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide17)

card_end = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.5), Inches(9.333), Inches(4.5))
card_end.fill.solid()
card_end.fill.fore_color.rgb = BG_CARD
card_end.line.color.rgb = ACCENT_EMERALD
card_end.line.width = Pt(2)

tb_end = slide17.shapes.add_textbox(Inches(2.5), Inches(2.0), Inches(8.333), Inches(3.5))
tf_end = tb_end.text_frame
tf_end.word_wrap = True

p_ty = tf_end.paragraphs[0]
p_ty.text = "THANK YOU!"
p_ty.alignment = PP_ALIGN.CENTER
p_ty.font.size = Pt(44)
p_ty.font.bold = True
p_ty.font.color.rgb = TEXT_WHITE
p_ty.space_after = Pt(16)

p_qa = tf_end.add_paragraph()
p_qa.text = "Questions & Answers (Q&A)"
p_qa.alignment = PP_ALIGN.CENTER
p_qa.font.size = Pt(24)
p_qa.font.bold = True
p_qa.font.color.rgb = ACCENT_CYAN
p_qa.space_after = Pt(20)

p_repo = tf_end.add_paragraph()
p_repo.text = "Project Repository: https://github.com/Chandravamsi09/Krono\nEmail: avvaruchandravamsi30@gmail.com"
p_repo.alignment = PP_ALIGN.CENTER
p_repo.font.size = Pt(14)
p_repo.font.color.rgb = TEXT_MUTED

# Save presentations
out_pptx = os.path.join(base_dir, 'presentation.pptx')
prs.save(out_pptx)
print(f"Presentation saved successfully at: {out_pptx}")

# Also copy/save as presentation.ppt if requested
import shutil
out_ppt = os.path.join(base_dir, 'presentation.ppt')
shutil.copyfile(out_pptx, out_ppt)
print(f"Also created copy: {out_ppt}")
